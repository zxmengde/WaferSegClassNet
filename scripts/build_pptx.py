#!/usr/bin/env python
"""
PPT文件生成脚本

使用python-pptx从SLIDES.md生成slides/final.pptx

用法:
    conda run -n wafer-seg-class python scripts/build_pptx.py --slides_md slides/SLIDES.md --results_root results --out slides/final.pptx

依赖:
    conda run -n wafer-seg-class pip install python-pptx
"""

import argparse
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple

PPTX_AVAILABLE = False
Presentation = None
Inches = None
Pt = None

def _check_pptx_import():
    """延迟检查python-pptx是否可用"""
    global PPTX_AVAILABLE, Presentation, Inches, Pt
    try:
        from pptx import Presentation as _Presentation
        from pptx.util import Inches as _Inches, Pt as _Pt
        Presentation = _Presentation
        Inches = _Inches
        Pt = _Pt
        PPTX_AVAILABLE = True
        return True
    except ImportError:
        PPTX_AVAILABLE = False
        return False


def parse_slides_md(filepath: Path) -> List[Dict]:
    """
    解析SLIDES.md文件，提取每页幻灯片的内容
    
    Returns:
        幻灯片列表，每个元素包含 {title, content, images}
    """
    if not filepath.exists():
        return []
    
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 按 "---" 分割幻灯片
    slides_raw = re.split(r'\n---\n', content)
    
    slides = []
    for slide_content in slides_raw:
        slide_content = slide_content.strip()
        if not slide_content:
            continue
        
        # 提取标题（## Slide N: xxx 或 ### xxx）
        title_match = re.search(r'^##\s+Slide\s+\d+:\s*(.+)$', slide_content, re.MULTILINE)
        if title_match:
            title = title_match.group(1).strip()
        else:
            # 尝试匹配 ### 标题
            title_match = re.search(r'^###\s+(.+)$', slide_content, re.MULTILINE)
            if title_match:
                title = title_match.group(1).strip()
            else:
                # 尝试匹配 # 标题
                title_match = re.search(r'^#\s+(.+)$', slide_content, re.MULTILINE)
                if title_match:
                    title = title_match.group(1).strip()
                else:
                    title = "Untitled"
        
        # 提取图片引用
        images = re.findall(r'!\[.*?\]\(([^)]+)\)', slide_content)
        
        # 清理内容（移除标题行和图片引用）
        clean_content = slide_content
        clean_content = re.sub(r'^##\s+Slide\s+\d+:.*$', '', clean_content, flags=re.MULTILINE)
        clean_content = re.sub(r'^###\s+.*$', '', clean_content, flags=re.MULTILINE)
        clean_content = re.sub(r'!\[.*?\]\([^)]+\)', '', clean_content)
        clean_content = clean_content.strip()
        
        slides.append({
            'title': title,
            'content': clean_content,
            'images': images
        })
    
    return slides


def add_title_slide(prs: 'Presentation', title: str, subtitle: str = "") -> None:
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle


def add_content_slide(prs: 'Presentation', title: str, content: str, 
                      images: List[str], results_root: Path) -> None:
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # 设置内容
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    
    # 解析内容为段落
    lines = content.split('\n')
    first_para = True
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # 跳过表格行
        if line.startswith('|'):
            continue
        
        # 处理列表项
        if line.startswith('- ') or line.startswith('* '):
            line = line[2:]
        elif re.match(r'^\d+\.\s', line):
            line = re.sub(r'^\d+\.\s', '', line)
        
        # 移除Markdown格式
        line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)  # 粗体
        line = re.sub(r'\*(.+?)\*', r'\1', line)  # 斜体
        line = re.sub(r'`(.+?)`', r'\1', line)  # 代码
        
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(18)
    
    # 添加图片（如果有且存在）
    for i, img_path in enumerate(images[:2]):  # 最多添加2张图片
        # 处理相对路径
        if img_path.startswith('../'):
            img_path = img_path[3:]
        
        full_path = results_root.parent / img_path
        if not full_path.exists():
            full_path = Path(img_path)
        
        if full_path.exists():
            try:
                left = Inches(5 + i * 2.5)
                top = Inches(2)
                width = Inches(2.5)
                slide.shapes.add_picture(str(full_path), left, top, width=width)
            except Exception as e:
                print(f"[WARN] 无法添加图片 {img_path}: {e}")



def build_pptx(slides_md: Path, results_root: Path, output_path: Path) -> bool:
    """
    从SLIDES.md构建PPTX文件
    
    Args:
        slides_md: SLIDES.md文件路径
        results_root: 实验结果根目录
        output_path: 输出PPTX文件路径
        
    Returns:
        是否成功生成
    """
    # 延迟导入检查
    if not _check_pptx_import():
        print("[ERROR] python-pptx未安装，无法生成PPTX")
        print("[INFO] 请运行: conda run -n wafer-seg-class pip install python-pptx")
        return False
    
    # 解析SLIDES.md
    slides = parse_slides_md(slides_md)
    
    if not slides:
        print("[ERROR] 未能从SLIDES.md解析出幻灯片内容")
        return False
    
    print(f"[INFO] 解析到 {len(slides)} 页幻灯片")
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    
    # 添加幻灯片
    for i, slide_data in enumerate(slides):
        title = slide_data['title']
        content = slide_data['content']
        images = slide_data['images']
        
        print(f"[INFO] 处理第 {i+1} 页: {title}")
        
        # 第一页使用标题布局
        if i == 0:
            # 提取副标题
            subtitle_match = re.search(r'\*\*副标题\*\*:\s*(.+)', content)
            subtitle = subtitle_match.group(1) if subtitle_match else ""
            add_title_slide(prs, title, subtitle)
        else:
            add_content_slide(prs, title, content, images, results_root)
    
    # 确保输出目录存在
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # 保存文件
    try:
        prs.save(str(output_path))
        print(f"[SUCCESS] PPTX已生成: {output_path}")
        return True
    except Exception as e:
        print(f"[ERROR] 保存PPTX失败: {e}")
        return False


def main():
    parser = argparse.ArgumentParser(
        description="从SLIDES.md生成PPTX文件",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
    conda run -n wafer-seg-class python scripts/build_pptx.py --slides_md slides/SLIDES.md --results_root results --out slides/final.pptx

依赖:
    conda run -n wafer-seg-class pip install python-pptx
        """
    )
    
    parser.add_argument(
        "--slides_md",
        type=str,
        default="slides/SLIDES.md",
        help="SLIDES.md文件路径 (默认: slides/SLIDES.md)"
    )
    
    parser.add_argument(
        "--results_root",
        type=str,
        default="results",
        help="实验结果根目录 (默认: results)"
    )
    
    parser.add_argument(
        "--out",
        type=str,
        default="slides/final.pptx",
        help="输出PPTX文件路径 (默认: slides/final.pptx)"
    )
    
    args = parser.parse_args()
    
    slides_md = Path(args.slides_md)
    results_root = Path(args.results_root)
    output_path = Path(args.out)
    
    if not slides_md.exists():
        print(f"[ERROR] SLIDES.md文件不存在: {slides_md}")
        print("[INFO] 请先运行: conda run -n wafer-seg-class python scripts/generate_slides_md.py")
        return 1
    
    if not results_root.exists():
        print(f"[WARN] 结果目录不存在: {results_root}")
        print("[INFO] 图片可能无法正确添加到PPTX")
    
    success = build_pptx(slides_md, results_root, output_path)
    
    return 0 if success else 1


if __name__ == "__main__":
    exit(main())
